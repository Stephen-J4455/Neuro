import os
import re
from huggingface_hub import InferenceClient
import requests
import fitz  # PyMuPDF
from pptx import Presentation

class NeuroAi:
    def __init__(self):
        self.patterns = [
            (re.compile(r"\b(who are you|who are u)\b", re.I),[
                "I am Neuro, a SMART ENTITY for LEARNING, MANAGEMENT and ASSISTANCE, a Ghanaian startup headquartered in The University of Professional Studies Accra."
            ]),
            (re.compile(r"\b(who created you|who created u)\b", re.I),[
                "I was developed by Stephen J. Amuzu, a Ghanaian software engineer and entrepreneur, as part of a project to create an AI assistant for learning and management."
            ])
        ]
        # Hugging Face client
        hf_token = os.environ.get("HF_TOKEN")
        self.hf_client = InferenceClient(token=hf_token) if hf_token else None

    def _download_file(self, url):
        try:
            response = requests.get(url, stream=True)
            response.raise_for_status()
            return response.content
        except requests.exceptions.RequestException as e:
            print(f"Error downloading file: {e}")
            return None

    def _extract_text_from_pdf(self, content):
        try:
            doc = fitz.open(stream=content, filetype="pdf")
            text = ""
            for page in doc:
                text += page.get_text()
            return text
        except Exception as e:
            print(f"Error extracting text from PDF: {e}")
            return None

    def _extract_text_from_pptx(self, content):
        try:
            from io import BytesIO
            prs = Presentation(BytesIO(content))
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            print(f"Error extracting text from PPTX: {e}")
            return None

    def ask(self, prompt, history=[], file_url=None):
        if not self.hf_client:
            return "Hugging Face API is not configured."

        file_content = ""
        if file_url:
            file_data = self._download_file(file_url)
            if file_data:
                if file_url.lower().endswith('.pdf'):
                    file_content = self._extract_text_from_pdf(file_data)
                elif file_url.lower().endswith(('.ppt', '.pptx')):
                    file_content = self._extract_text_from_pptx(file_data)

        if file_content:
            prompt = f"Use the following content to answer the question: {file_content}\n\nQuestion: {prompt}"

        messages = []
        for message in history:
            if message.get('sender') == 'user':
                messages.append({'role': 'user', 'content': message.get('text')})
            elif message.get('sender') == 'ai':
                messages.append({'role': 'assistant', 'content': message.get('text')})

        messages.append({"role": "user", "content": prompt})

        try:
            response = self.hf_client.chat.completions.create(
                model="Qwen/Qwen3-Next-80B-A3B-Instruct",
                messages=messages,
                max_tokens=2048 # Lowered to reduce memory usage
            )
            content = response.choices[0].message.content
            print(f"[DEBUG] Hugging Face response length: {len(content) if content else 0}")
            if content and len(content) < 3500:
                print("[DEBUG] Response may be incomplete or truncated:", content)
            return content
        except Exception as e:
            print(f"[DEBUG] Hugging Face API error: {e}")
            return f"Error with Hugging Face API: {e}"

    def chat(self):
        print("SELMA: Hi! I'm your personal chat assistant. What would you like to do?")
        while True:
            try:
                user_input = input("You: ").strip()
                if not user_input:
                    continue
                response = self.ask(user_input)
                print(f"SELMA: {response}")
            except KeyboardInterrupt:
                print("\nSELMA: Goodbye!")
                break
            except Exception as e:
                print(f"SELMA: Oops! Error: {e}")


if __name__ == "__main__":
    bot = NeuroAi()
    bot.chat()
